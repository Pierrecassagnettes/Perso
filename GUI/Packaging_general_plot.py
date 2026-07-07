# -*- coding: utf-8 -*-
"""
Created on Thu May 27 09:26:54 2026

@author: Pierre.Cassagnettes

Goal of this graphical interface : For a given Tapeout name, this script extracts metrics to plot the exact number of occurrences for
        your chosen 'Final Binning'. Wafers are automatically discovered via the network (W: drive).
        All output figures are automatically exported and saved directly into your chosen destination folder below.

"""

import os
import sys
from pathlib import Path
import yaml
import textwrap
import numpy as np
import tempfile

import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.cm as cm
import matplotlib
from matplotlib.lines import Line2D

matplotlib.use("Agg")
import matplotlib.image as mpimg

import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import threading
from PIL import Image, ImageTk
import WaferMap

# Ordered list of Final Binning options from the provided table (Best to Worst)
BINNING_OPTIONS = [
    "1 - OK (best)",
    "13 - mpd short",
    "14 - mpd open",
    "17 - temp sensor open",
    "16 - ring open",
    "23 - ring resistance",
    "40 - deembedded power",
    "41 - comb length",
    "42 - max abs f i sigma",
    "43 - min deembedded power",
    "44 - range deembedded power",
    "30 - trim fit period",
    "15 - mzi open",
    "22 - mzi resistance",
    "20 - laser resistance",
    "11 - laser open",
    "12 - laser short (worst)"
]

def display_image(root, file_path):
    """Display the image"""
    if not file_path or not os.path.exists(file_path):
        return
        
    img_window = tk.Toplevel(root)
    img_window.title(f"Preview: {os.path.basename(file_path)}")

    img = Image.open(file_path)
    
    screen_w = root.winfo_screenwidth()
    screen_h = root.winfo_screenheight()
    
    max_w = int(screen_w * 0.75)
    max_h = int(screen_h * 0.75)
    
    img.thumbnail((max_w, max_h), Image.Resampling.LANCZOS)
    photo = ImageTk.PhotoImage(img)

    label = tk.Label(img_window, image=photo)
    label.image = photo 
    label.pack(padx=10, pady=10)


def filter_wafers(filters, wafers):
    """find wafers that matches all filters"""
    result = {}
    for k, w in wafers.items():
        if all(str(w.get(f)) == v for f, v in filters.items() if v):
            result[k] = w
    return result


def find_files(root_dir, wafer_list, base_name, tapeout):
    """find all files related to wafers of wafer_list
    if it is Tapeout 5 : with the excel report (generated with Automatisation_packaging.py )
    if it is an other Tapeout : with raw datas files
    """
    files = []
    wafers = []
    
    for wafer in wafer_list:
        for temp, pr in [("50C", "/processed_data"), ("50deg", "/processed data")]:
            root_path = os.path.join(root_dir, wafer, "Leaflight", temp, base_name + pr)
            
            if tapeout == "TO5":
                target = os.path.join(root_path, f"report_{wafer}_{base_name}.xlsx")
                if os.path.exists(target):
                    files.append((wafer, target))
                    wafers.append(wafer)
                    break
            else:
                parent_path = os.path.join(root_dir, wafer, "Leaflight", temp, base_name)
                raw_data_path = os.path.join(parent_path, "raw_data")
                if not os.path.exists(raw_data_path):
                    raw_data_path = os.path.join(parent_path, "raw data")
                
                if os.path.exists(root_path) or os.path.exists(raw_data_path):
                    files.append((wafer, root_path))
                    wafers.append(wafer)
                    break
                    
    return files, wafers


def prepare_boxplot_data(df, group_col="Wafer name", value_col="Result"):
    """prepare datas for the boxplot"""
    grouped = df.groupby(group_col, observed=True)
    data = []
    labels = []

    for name, group in grouped:
        values = pd.to_numeric(group[value_col], errors="coerce").dropna()
        if len(values) == 0:
            continue 
        data.append(values.values)
        labels.append(str(name))

    return data, labels


def plot_boxplot_with_points(
    data,
    labels,
    ylabel,
    title=None,
    jitter=0.08,
    logscale=False,
    spec=None,
    figsize=(8, 6),
    fontsize_labels=12,
    fontsize_ticks=10,
    fontsize_title=14,
    ymin=None,
    ymax=None,
    save_path=None
):
    """plot the boxplot with given spec ( float or range of floats ), ymin ymax (limits of y axis)"""
    fig, ax = plt.subplots(figsize=figsize)
    positions = np.arange(1, len(data) + 1)

    bp = ax.boxplot(
        data,
        positions=positions,
        labels=labels,
        showfliers=False,
        patch_artist=True
    )

    for box in bp["boxes"]:
        box.set(facecolor="lightgray", alpha=0.7)

    for pos, values in zip(positions, data):
        x = np.random.normal(pos, jitter, size=len(values))
        ax.scatter(x, values, s=30, alpha=0.8)

    if spec is not None:
        if isinstance(spec, (int, float)):
            ax.axhline(spec, color="red", linestyle="--", linewidth=2, label="Spec")
        elif isinstance(spec, (tuple, list)) and len(spec) == 2:
            ax.axhspan(spec[0], spec[1], color="red", alpha=0.15, label="Spec range")
        elif isinstance(spec, dict):
            for pos, label in zip(positions, labels):
                if label not in spec:
                    continue
                wafer_spec = spec[label]
                if isinstance(wafer_spec, (int, float)):
                    ax.hlines(wafer_spec, pos - 0.35, pos + 0.35, colors="red", linestyles="--", linewidth=2, label='Spec')
                elif isinstance(wafer_spec, (tuple, list)) and len(wafer_spec) == 2:
                    ax.fill_between([pos - 0.35, pos + 0.35], wafer_spec[0], wafer_spec[1], color="red", alpha=0.2, label='Spec range')
                else:
                    raise ValueError(f"Spec invalide pour le wafer {label}")
        else:
            raise ValueError("spec doit être None, float, tuple(min,max) ou dict")

    ax.set_ylabel(ylabel, fontsize=fontsize_labels)
    ax.tick_params(axis="x", labelsize=fontsize_ticks, rotation=45)
    ax.tick_params(axis="y", labelsize=fontsize_ticks)

    if spec:
        ax.legend(fontsize=fontsize_ticks)
    
    if title:
        wrapped_title = "\n".join(textwrap.wrap(title, width=50))
        ax.set_title(wrapped_title, fontsize=fontsize_title, pad=10)

    if logscale:
        ax.set_yscale("log")

    ax.grid(axis="y", linestyle="--", alpha=0.5)
    ax.grid(axis="x", linestyle="--", alpha=0.5)
    
    if ymin is not None or ymax is not None:
        ax.set_ylim(
            bottom=ymin if ymin is not None else ax.get_ylim()[0],
            top=ymax if ymax is not None else ax.get_ylim()[1],
        )

    plt.tight_layout()
    if save_path:
        plt.savefig(save_path, dpi=300)
        plt.close(fig)
        return save_path
    return fig


def plot_final_binning(files, final_binning, saving_path, channel, base_name, tapeout, screen_w, screen_h, plot_choices, save_plots):
    """main function that plots all figures"""
    wafers = []
    counts = []
    
    dies_labels = []
    dies_min = []
    dies_max = []
    dies_wafers = []
    dies_wavelength = []
    
    bin_desc = next((opt for opt in BINNING_OPTIONS if opt.startswith(f"{final_binning} - ")), str(final_binning))
    
    dynamic_width = max(14, (screen_w * 0.009))
    dynamic_height = max(8, (screen_h * 0.009))
    figsize = (dynamic_width, dynamic_height)
    
    for wafer, file in files:
        if tapeout == "TO5":
            df = pd.read_excel(file)
            required_cols = {"X", "Y", "Final Binning", "expected power output (dBm)", "Channel", "wavelength (nm)", "Lasers/heater"}
            if not required_cols.issubset(df.columns):
                raise ValueError(f"The file {file} does not contain the required columns: {required_cols}")

            df["expected power output (dBm)"] = pd.to_numeric(df["expected power output (dBm)"], errors="coerce")
            df["wavelength (nm)"] = pd.to_numeric(df["wavelength (nm)"], errors="coerce")
            df = df.dropna(subset=["X", "Y"])
            
            df_unique = df.drop_duplicates(subset=["X", "Y"])
            count = len(df_unique[df_unique["Final Binning"] == int(final_binning)])

            df_filtered = df[(df["Final Binning"] == int(final_binning)) & (df["Channel"] == channel)]
            grouped_power = df_filtered.groupby(["X", "Y"])["expected power output (dBm)"].agg(["min", "max"]).reset_index()
            
            df_wave = df_filtered[df_filtered["Lasers/heater"].isin(["F0", "F1"])]
            grouped_wave = df_wave.groupby(["X", "Y"])["wavelength (nm)"].first().reset_index()
            grouped_all = pd.merge(grouped_power, grouped_wave, on=["X", "Y"], how="left")

            for _, row in grouped_all.iterrows():
                dies_labels.append(f"{wafer}\n({int(row['X'])},{int(row['Y'])})")
                dies_min.append(row["min"])
                dies_max.append(row["max"])
                dies_wafers.append(wafer)
                dies_wavelength.append(row["wavelength (nm)"])
                
        else:
            file_dir = file
            parent_dir = os.path.dirname(file_dir)
            raw_data_dir = os.path.join(parent_dir, "raw_data")
            if not os.path.exists(raw_data_dir):
                raw_data_dir = os.path.join(parent_dir, "raw data")
            
            parquet_path = os.path.join(file_dir, "metrics", "flattened_metrics.parquet")
            if not os.path.exists(parquet_path):
                raise FileNotFoundError(f"Le fichier de métriques est introuvable : {parquet_path}")
            
            df_metrics = pd.read_parquet(parquet_path)
            df_metrics = df_metrics.rename(columns={"device_x": "X", "device_y": "Y"})
            
            df_metrics = df_metrics[df_metrics["metric_name"] == "deembedded_output_power"]
            df_metrics["laser_num"] = pd.to_numeric(df_metrics["meta_id"].str[1:], errors="coerce")
            df_metrics = df_metrics[df_metrics["laser_num"] % 2 == 0]
            
            binning_dir = os.path.join(file_dir, "binning")
            target_bin_dir = binning_dir if os.path.exists(binning_dir) else file_dir
            
            binning_file = os.path.join(target_bin_dir, "binning.csv")
            if not os.path.exists(binning_file):
                binning_file = os.path.join(target_bin_dir, "binning_report.csv")
            if not os.path.exists(binning_file):
                try:
                    matches = [f for f in os.listdir(target_bin_dir) if "binning" in f.lower() and f.endswith(".csv")]
                    if matches:
                        binning_file = os.path.join(target_bin_dir, matches[0])
                    else:
                        raise FileNotFoundError()
                except Exception:
                    raise FileNotFoundError(f"Le fichier de binning (.csv) est introuvable : {target_bin_dir}")
            
            df_binning = pd.read_csv(binning_file)
            df_binning = df_binning.rename(columns={"Final Bin": "Final Binning"})
            df_binning["Final Binning"] = pd.to_numeric(df_binning["Final Binning"], errors="coerce")
            df_binning = df_binning.dropna(subset=["X", "Y"])
            
            df_unique = df_binning.drop_duplicates(subset=["X", "Y"])
            count = len(df_unique[df_unique["Final Binning"] == int(final_binning)])
            
            df_bin_filtered = df_binning[df_binning["Final Binning"] == int(final_binning)]
            df_bin_filtered = df_bin_filtered.drop_duplicates(subset=["X", "Y"])
            
            try:
                raw_files = [f for f in os.listdir(raw_data_dir) if f.endswith(".txt") and "_spectrum_" in f]
            except Exception:
                raw_files = []
                
            for _, bin_row in df_bin_filtered.iterrows():
                bx, by = int(bin_row["X"]), int(bin_row["Y"])
                
                df_die_metrics = df_metrics[(df_metrics["X"] == bx) & (df_metrics["Y"] == by)]
                die_powers = pd.to_numeric(df_die_metrics["y"], errors="coerce").dropna().tolist()
                f0_wavelength = float('nan')
                
                if tapeout == "TO4.1":
                    for rf in raw_files:
                        parts = rf.split("_")
                        if len(parts) > 5 and parts[0] == "Die" and parts[1] == str(bx) and parts[2] == str(by) and parts[5] == "F0":
                            try:
                                filepath = os.path.join(raw_data_dir, rf)
                                df_txt = pd.read_csv(filepath, sep="\t")
                                if "Power__dBm" in df_txt.columns and "Wavelength__nm" in df_txt.columns:
                                    idx_max = pd.to_numeric(df_txt["Power__dBm"], errors="coerce").idxmax()
                                    if pd.notna(idx_max):
                                        f0_wavelength = pd.to_numeric(df_txt.loc[idx_max, "Wavelength__nm"], errors="coerce")
                                break
                            except Exception:
                                pass
                            
                if die_powers:
                    dies_labels.append(f"{wafer}\n({bx},{by})")
                    dies_min.append(min(die_powers) + 21)
                    dies_max.append(max(die_powers) + 21)
                    dies_wafers.append(wafer)
                    dies_wavelength.append(f0_wavelength)

        wafers.append(wafer)
        counts.append(count)

    save1, save2, save3 = None, None, None
    suffix_chan = f"_{channel}" if tapeout == "TO5" else "_all"
    target_dir = saving_path if save_plots else tempfile.gettempdir()

    # ---- Plot 1: Total Dies Count ----
    if plot_choices.get("plot1", True):
        fig, ax = plt.subplots(figsize=figsize)
        bars = ax.bar(wafers, counts)
        ax.set_xlabel("Wafer")
        ax.set_ylabel(f"Number of Dies with Final Binning = {final_binning}")
        ax.set_title(f"Number of Dies with Final Binning = {bin_desc}")
        plt.xticks(rotation=0)
        plt.grid(axis='y')
        
        max_y = max(counts) if max(counts) > 0 else 1
        ax.set_ylim(0, max_y * 1.2)
        
        symbol_added_p1 = False
        for bar, c in zip(bars, counts):
            if c == 0:
                x_pos = bar.get_x() + bar.get_width()/2
                ax.text(x_pos, max_y * 0.1, r'$\oslash$', fontsize=28, color='#1d5073', ha='center', va='center', weight='bold')
                symbol_added_p1 = True
                
        if symbol_added_p1:
            
            lbl_text = "No fully functional LEAF Light" if int(final_binning) == 1 else f"No final binning {final_binning}"
            legend_elements = [Line2D([0], [0], marker=r'$\oslash$', color='none', markeredgecolor='#1d5073', markersize=14, label=lbl_text)]
            ax.legend(handles=legend_elements, loc='upper right')
            
        plt.tight_layout()
        save1 = os.path.join(target_dir, f"General_plot_{base_name}_{tapeout}{suffix_chan}.png")
        plt.savefig(save1, dpi=300)
        plt.close(fig)

    # ---- Plot 2: Min/Max Power Range per Die ----
    if plot_choices.get("plot2", True):
        fig, ax = plt.subplots(figsize=figsize)
        x_positions = []
        x_labels = []
        current_pos = 0
        spacing_between_wafers = 2.5
        spacing_between_dies = 0.4
        
        colors = matplotlib.colormaps["tab10"]
        empty_wafer_positions_p2 = []

        for i, wafer in enumerate(wafers):
            if i > 0:
                ax.axvline(x=current_pos - (spacing_between_wafers / 2), color='gray', linestyle=':', alpha=0.6)
                
            wafer_indices = [idx for idx, w in enumerate(dies_wafers) if w == wafer]
            color = colors(i % 10)
            
            if wafer_indices:
                for j, idx in enumerate(wafer_indices):
                    x = current_pos + j * spacing_between_dies
                    ax.plot([x, x], [dies_min[idx], dies_max[idx]], marker='o', color=color)
                    label_text = dies_labels[idx].split("\n")[1]
                    ax.text(x, dies_max[idx] + 0.1 , label_text, rotation=90, fontsize=8, ha='center', va='bottom', color=color)

                center = current_pos + (len(wafer_indices) - 1) * spacing_between_dies / 2
                x_positions.append(center)
                current_pos += len(wafer_indices) * spacing_between_dies + spacing_between_wafers
            else:
                empty_wafer_positions_p2.append(current_pos)
                x_positions.append(current_pos)
                current_pos += spacing_between_wafers
                
            x_labels.append(wafer)

        ax.set_xticks(x_positions)
        ax.set_xticklabels(x_labels)
        ax.set_xlabel("Wafer")
        ax.set_ylabel("Power (dBm)")
        title_chan = f", Channel={channel}" if tapeout == "TO5" else ""
        
        if tapeout == "TO5":
            ax.set_title(f"Min/Max Power (on chip after mux) per Die (Final Binning={bin_desc}{title_chan}) ")
        else:
            ax.set_title(f"Min/Max Power (on chip before mux) per Die (Final Binning={bin_desc}{title_chan})\n only considering even lasers")
            
        y_min, y_max = ax.get_ylim()
        y_center = (y_min + y_max) / 2
        for pos in empty_wafer_positions_p2:
            ax.text(pos, y_center, r'$\oslash$', fontsize=28, color='#1d5073', ha='center', va='center', weight='bold')
            
        if empty_wafer_positions_p2:
            
            lbl_text = "No fully functional LEAF Light" if int(final_binning) == 1 else f"No final binning {final_binning}"
            legend_elements = [Line2D([0], [0], marker=r'$\oslash$', color='none', markeredgecolor='#1d5073', markersize=14, label=lbl_text)]
            ax.legend(handles=legend_elements, loc='upper right')

        plt.grid(axis='y')
        plt.tight_layout()
        save2 = os.path.join(target_dir, f"Power_per_die_{base_name}_{tapeout}{suffix_chan}.png")
        plt.savefig(save2, dpi=300)
        plt.close(fig)
    
    # ---- Plot 3: Laser Wavelength  ----
    if plot_choices.get("plot3", True) and tapeout in ["TO5", "TO4.1"]:
        fig, ax = plt.subplots(figsize=figsize)
        x_positions = []
        x_labels = []
        current_pos = 0
        spacing_between_wafers = 2.5
        spacing_between_dies = 0.4
        colors = matplotlib.colormaps["tab10"]
        empty_wafer_positions_p3 = []

        for i, wafer in enumerate(wafers):
            if i > 0:
                ax.axvline(x=current_pos - (spacing_between_wafers / 2), color='gray', linestyle=':', alpha=0.6)
                
            wafer_indices = [idx for idx, w in enumerate(dies_wafers) if w == wafer]
            color = colors(i % 10)
            
            if wafer_indices:
                for j, idx in enumerate(wafer_indices):
                    x = current_pos + j * spacing_between_dies
                    wave_val = dies_wavelength[idx]
                    if not pd.isna(wave_val):
                        ax.plot(x, wave_val, marker='o', color=color)
                        label_text = dies_labels[idx].split("\n")[1]
                        ax.text(x, wave_val + 0.1, label_text, rotation=90, fontsize=8, ha='center', va='bottom', color=color)

                center = current_pos + (len(wafer_indices) - 1) * spacing_between_dies / 2
                x_positions.append(center)
                current_pos += len(wafer_indices) * spacing_between_dies + spacing_between_wafers
            else:
                empty_wafer_positions_p3.append(current_pos)
                x_positions.append(current_pos)
                current_pos += spacing_between_wafers
                
            x_labels.append(wafer)

        ax.set_xticks(x_positions)
        ax.set_xticklabels(x_labels)
        ax.set_xlabel("Wafer")
        ax.set_ylabel("Wavelength (nm)")
        title_chan = f", Channel={channel}" if tapeout == "TO5" else ""
        ax.set_title(f"Wavelength F0 (nm) (Final Binning={bin_desc}{title_chan if tapeout=='TO5' else ''})")
        
        y_min_w, y_max_w = ax.get_ylim()
        y_center_w = (y_min_w + y_max_w) / 2 if y_max_w != 1 else 0
        for pos in empty_wafer_positions_p3:
            ax.text(pos, y_center_w, r'$\oslash$', fontsize=28, color='#1d5073', ha='center', va='center', weight='bold')
            
        if empty_wafer_positions_p3:
            
            lbl_text = "No fully functional LEAF Light" if int(final_binning) == 1 else f"No final binning {final_binning}"
            legend_elements = [Line2D([0], [0], marker=r'$\oslash$', color='none', markeredgecolor='#1d5073', markersize=14, label=lbl_text)]
            ax.legend(handles=legend_elements, loc='upper right')

        plt.grid(axis='y')
        ax.ticklabel_format(style='plain', axis='y')
        plt.tight_layout()

        save3 = os.path.join(target_dir, f"wavelength_{base_name}_{tapeout}_{channel if tapeout=='TO5' else 'all'}.png")
        plt.savefig(save3, dpi=300)
        plt.close(fig)
    
    return save1, save2, save3


def plot_die_power_range(file_path, channel, selected_wafer, tapeout):
    """Min/Max Power Range per Die for all dies when a wafer is selected """
    if tapeout == "TO5":
        df = pd.read_excel(file_path, engine="openpyxl")
        df["expected power output (dBm)"] = pd.to_numeric(df["expected power output (dBm)"], errors="coerce")
        df = df.dropna(subset=["expected power output (dBm)", "X", "Y"])
        df = df[(df["Channel"] == channel)]
        
        # Identification des lasers en panne avant filtrage (< -50 dBm)
        df["is_dead"] = df["expected power output (dBm)"] < -50.0
        dead_lasers = df.groupby(["X", "Y"])["is_dead"].sum().reset_index()
        
        # Filtrage des puissances valides
        df_clean = df[df["expected power output (dBm)"] >= -50.0]
        grouped = df_clean.groupby(["X", "Y"])["expected power output (dBm)"].agg(["min", "max"]).reset_index()
        grouped_max = df_clean.groupby(["X", "Y"])["expected power output (dBm)"].max().reset_index()
        grouped_max.columns = ["X", "Y", "expected power output (dBm)"]
        
        # Re-fusionner le compte des lasers morts
        grouped = pd.merge(grouped, dead_lasers, on=["X", "Y"], how="left")
        grouped_max = pd.merge(grouped_max, dead_lasers, on=["X", "Y"], how="left")
    else:
        file_dir = file_path
        parquet_path = os.path.join(file_dir, "metrics", "flattened_metrics.parquet")
        
        if os.path.exists(parquet_path):
            df_all = pd.read_parquet(parquet_path)
            df_all = df_all.rename(columns={"device_x": "X", "device_y": "Y", "y": "power"})
            
            df_all = df_all[df_all["metric_name"] == "deembedded_output_power"]
            df_all["laser_num"] = pd.to_numeric(df_all["meta_id"].str[1:], errors="coerce")
            df_all = df_all[df_all["laser_num"] % 2 == 0]
            
            df_all["power"] = pd.to_numeric(df_all["power"], errors="coerce")
            df_all["power"] = df_all["power"] +21
            df_all = df_all.dropna(subset=["X", "Y"])
            limite_power = -10
            # Identification des lasers en panne avant filtrage (< -50 dBm)
            df_all["is_dead"] = df_all["power"] < limite_power
            dead_lasers = df_all.groupby(["X", "Y"])["is_dead"].sum().reset_index()
            
            # Filtrage des puissances valides
            df_clean = df_all[df_all["power"] >= limite_power]
            df_clean = df_clean[df_clean["power"] <= 25]
            grouped = df_clean.groupby(["X", "Y"])["power"].agg(["min", "max"]).reset_index()
            grouped_max = df_clean.groupby(["X", "Y"])["power"].max().reset_index()
            grouped_max.columns = ["X", "Y", "expected power output (dBm)"]
            
            grouped = pd.merge(grouped, dead_lasers, on=["X", "Y"], how="left")
            grouped_max = pd.merge(grouped_max, dead_lasers, on=["X", "Y"], how="left")
        else:
            grouped = pd.DataFrame(columns=["X", "Y", "min", "max", "is_dead"])
            grouped_max = pd.DataFrame(columns=["X", "Y", "expected power output (dBm)", "is_dead"])
        
    grouped_max["Die"] = grouped_max.apply(lambda row: f"({int(row['X'])},{int(row['Y'])})", axis=1)
    
    base_path = r"W:\50-DEVELOPMENT\TEST\Temporary data database\Data"
    map_position = os.path.join(base_path, selected_wafer, 'Map infos')
    
    if not os.path.isfile(os.path.join(map_position, "Die_vs_XY.txt")):
        map_position = os.path.join(base_path, 'D24S2097_23', 'Map infos')
        
    with open(os.path.join(map_position, "Die_vs_XY.txt"), "r") as f:
        lines = f.readlines()

    xy_to_die = {}
    center_x, center_y = 4, 4
        
    for row_idx, line in enumerate(lines):
        values = line.strip().split("\t")
        for col_idx, value in enumerate(values):
            die_number = int(value)
            if die_number != 0:
                x = col_idx - center_x
                y = row_idx - center_y
                xy_to_die[(x, y)] = die_number

    grouped = grouped.sort_values(by=["X", "Y"])
    grouped["Die"] = grouped.apply(lambda row: f"({int(row['X'])},{int(row['Y'])})", axis=1)
    
    # ---- 1. GÉNERATION DU GRAPH FLATTENED WAFERMAP ----
    fig1, ax1 = plt.subplots(figsize=(12, 6))
    data_dict = {}
    for _, row in grouped_max.iterrows():
        x, y = int(row["X"]), int(row["Y"])
        max_power = row["expected power output (dBm)"]
        die_number = xy_to_die.get((x, y))
        if die_number is not None and pd.notna(max_power):
            data_dict[f"die_{die_number}"] = max_power

    title = f"WaferMap {selected_wafer} - Max Expected Power Output (dBm)"
    try:
        WaferMap.main(data_dict, fig1, ax1, title, wafer_name=selected_wafer)
    except Exception as e:
        print(f"Error drawing WaferMap module layout: {e}")
    plt.tight_layout()
    plt.close(fig1)

    # ---- 2. GÉNERATION DU GRAPH MIN-MAX POWER PAR DIE ----
    fig2, ax2 = plt.subplots(figsize=(14, 8))
    
    positions = np.arange(len(grouped))
    labels = grouped["Die"].tolist()
    
    # Tracé des barres verticales min/max avec ronds aux extrémités (Style Plot 2)
    for idx, row in grouped.iterrows():
        pos = positions[idx]
        
        # On ne trace la barre que si des données valides existent (>= -50 dBm)
        if pd.notna(row["min"]) and pd.notna(row["max"]):
            # Ligne verticale reliant le min au max
            ax2.plot([pos, pos], [row["min"], row["max"]], marker='o', color='#1f77b4', linewidth=1.5)
            
            # Étiquette de coordonnées (X,Y) positionnée juste au-dessus du max (+ 0.5 dBm)
            ax2.text(pos, row["max"] + 0.5, row["Die"], rotation=90, fontsize=8, ha='center', va='bottom', color='#1f77b4')
            
            # Si des lasers ont été coupés/filtrés car inférieurs à -50 dBm
            if row["is_dead"] > 0:
                msg = f"{int(row['is_dead'])} HS" if row["is_dead"] == 1 else f"{int(row['is_dead'])} HS"
                # Annotation textuelle placée encore un peu plus haut au-dessus de l'étiquette de coordonnée
                ax2.text(pos, row["min"] -0.6, msg, rotation=0, fontsize=7, ha='center', va='bottom', color='red', weight='bold')
                
        else:
            if row["is_dead"] > 0:
                ax2.text(pos, 0.0, r'$\oslash$', fontsize=18, color='red', ha='center', va='center')
                ax2.text(pos, 2.0, "All lasers HS", rotation=90, fontsize=7, ha='center', va='bottom', color='red', weight='bold')  
                         
    
    lbl_text = "HS = Laser's power on chip before mux <= -10 dBm"
    legend_elements = [Line2D([0], [0], color='none', markeredgecolor='#1d5073', markersize=14, label=lbl_text)]
    ax2.legend(handles=legend_elements, loc='upper right')
    ax2.set_xticks(positions)
    ax2.set_xticklabels(labels, rotation=90)
    ax2.set_xlabel("Die (X, Y)")
    ax2.set_ylabel("Expected Power Output (dBm)" if tapeout == "TO5" else "Deembedded Power (dBm)")
    ax2.set_title(f"Min-Max Range per Die | Wafer: {selected_wafer}")
    ax2.grid(axis='y', linestyle='--', alpha=0.7)
    fig2.tight_layout()

    return fig2, ax2


def generate_pptx_report(tapeout, final_binning, bin_desc, saving_path, pptx_records, save4=None):
    """generate ppt report with all selected figures for all bases avialable for the tapeout selected """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("python-pptx package missing.")
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # --- SLIDE 1 : Titre Page de Garde ---
    slide1 = prs.slides.add_slide(blank_layout)
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    
    if int(final_binning) == 1:
        p1.text = f"Summary of {tapeout} wafers' test\n- fully functional LEAF Light\n(bin 1)"
    else:
        p1.text = f"Summary of {tapeout} wafers' test\n- {bin_desc}"
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.name = "Arial"
    p1.alignment = PP_ALIGN.CENTER

    footer_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(2.2))
    tf2 = footer_box.text_frame
    tf2.word_wrap = True
    
    
    p3 = tf2.add_paragraph()
    p3.text = "June 2026."
    p3.font.size = Pt(16)
    p3.font.name = "Arial"
    p3.alignment = PP_ALIGN.CENTER
    
    p4 = tf2.add_paragraph()
    p4.text = "Scintil confidential"
    p4.font.size = Pt(16)
    p4.font.name = "Arial"
    p4.alignment = PP_ALIGN.CENTER

    if save4 and os.path.exists(save4):
        slide_ref = prs.slides.add_slide(blank_layout)
        tx_box = slide_ref.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
        p_ref = tx_box.text_frame.paragraphs[0]
        p_ref.text = "Reference Power Analysis - Overview of Plotted Metrics"
        p_ref.font.size = Pt(24)
        p_ref.font.bold = True
        p_ref.alignment = PP_ALIGN.CENTER
        
        with Image.open(save4) as im:
            r = im.size[0] / im.size[1]
        w, h = (11.5, 11.5 / r) if r > (11.5 / 5.5) else (5.5 * r, 5.5)
        slide_ref.shapes.add_picture(save4, Inches((13.333 - w)/2), Inches(1.3 + (5.5 - h)/2), width=Inches(w), height=Inches(h))

    def add_scaled_img(slide, img_path):
        """puts images with correct scale in slides"""
        with Image.open(img_path) as im:
            r = im.size[0] / im.size[1]
        w, h = (11.5, 11.5 / r) if r > (11.5 / 6.0) else (6.0 * r, 6.0)
        slide.shapes.add_picture(img_path, Inches((13.333 - w)/2), Inches(0.6 + (6.3 - h)/2), width=Inches(w), height=Inches(h))

    for record in pptx_records:
        base = record["base"]
        img1 = record.get("plot1")
        img2 = record.get("plot2")
        img3 = record.get("plot3")

        slide_div = prs.slides.add_slide(blank_layout)
        div_box = slide_div.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.333), Inches(1))
        p_div = div_box.text_frame.paragraphs[0]
        p_div.text = f"Bin {final_binning} - {base}"
        p_div.font.size = Pt(36)
        p_div.font.bold = True
        p_div.alignment = PP_ALIGN.CENTER

        if img1 and os.path.exists(img1):
            slide_p1 = prs.slides.add_slide(blank_layout)
            add_scaled_img(slide_p1, img1)

        if img2 and os.path.exists(img2):
            slide_p2 = prs.slides.add_slide(blank_layout)
            add_scaled_img(slide_p2, img2)

        if img3 and os.path.exists(img3):
            slide_p3 = prs.slides.add_slide(blank_layout)
            add_scaled_img(slide_p3, img3)

    output_path = os.path.join(saving_path, f"Summary_Report_{tapeout}_Bin_{final_binning}.pptx")
    prs.save(output_path)
    return output_path


def main(base_name, final_binning, channel, tapeout, saving_path, wafer_list, selected_wafer, generate_pptx, screen_w, screen_h, plot_choices, save_plots):
    root_data = r"W:\50-DEVELOPMENT\TEST\Temporary data database\Data"
    
    if tapeout in ["TO4", "TO4.1"]:
        bases_to_process = ["LA1"]
    else:
        bases_to_process = ["BASE_LONG", "BASE_SHORT", "ENH_LONG", "ENH_LONG_2x8"]
        
    pptx_records = []
    sel_save1, sel_save2, sel_save3, sel_save4 = None, None, None, None
    
    loop_bases = bases_to_process if generate_pptx else [base_name]
    dynamic_figsize = (max(14, (screen_w * 0.009)), max(8, (screen_h * 0.009)))
    
    for b in loop_bases:
        files = find_files(root_data, wafer_list, b, tapeout)[0]
        s1, s2, s3 = plot_final_binning(files, final_binning, saving_path, channel, b, tapeout, screen_w, screen_h, plot_choices, save_plots)
        
    
        
        s4 = None
        if b == base_name and selected_wafer:
            path_selected_wafer = next((f for w, f in files if w == selected_wafer), None)
            if path_selected_wafer:
                fig, ax = plot_die_power_range(path_selected_wafer, channel, selected_wafer, tapeout)
                ax.set_xlabel("Die")
                ax.set_ylabel("Min, Max Expected Power Output (dBm)")
                suffix_chan = f"_{channel}" if tapeout == "TO5" else "_all"
                ax.set_title(f"Min, Max Expected Power Output (dBm) | Base: {b} | Wafer: {selected_wafer}")
                
                title = f"Min_Max_expected_power_output_(dBm)_{b}_{selected_wafer}{suffix_chan}.png"
                target_dir = saving_path if save_plots else tempfile.gettempdir()
                s4 = os.path.join(target_dir, title)
                plt.savefig(s4, dpi=300)
                plt.close(fig)
                
        if b == base_name:
            sel_save1, sel_save2, sel_save3, sel_save4 = s1, s2, s3, s4
            
        pptx_records.append({
            "base": b,
            "plot1": s1,
            "plot2": s2,
            "plot3": s3
        })
        
    if generate_pptx:
        bin_desc = next((opt for opt in BINNING_OPTIONS if opt.startswith(f"{final_binning} - ")), str(final_binning))
        generate_pptx_report(tapeout, final_binning, bin_desc, saving_path, pptx_records, sel_save4)
        
    return sel_save1, sel_save2, sel_save3, sel_save4


def request_parameters():
    """graphical interface (GUI) part """
    def choose_folder():
        folder = filedialog.askdirectory()
        if folder:
            entry_folder.delete(0, tk.END)
            entry_folder.insert(0, folder)
    
    def start_processing(base_name, final_binning, channel, tapeout, saving_file, wafer_list, selected_wafer, generate_pptx, screen_w, screen_h, plot_choices, save_plots, display_plots):
        def run():
            try:
                save1, save2, save3, save4 = main(
                    base_name, final_binning, channel, tapeout, saving_file, wafer_list, selected_wafer, 
                    generate_pptx, screen_w, screen_h, plot_choices, save_plots
                )
                root.after(0, lambda: finish_processing(save1, save2, save3, save4, display_plots))
            except Exception as e:
                root.after(0, lambda: error_handler(e))

        threading.Thread(target=run, daemon=True).start()

    def error_handler(error):
        progress.stop()
        progress.grid_forget()
        btn_run.config(state="normal")
        messagebox.showerror("Critical Error", f"An error occurred during automated processing:\n{error}")

    def choose_wafer():
        base_name = entry_base.get().strip()
        combobox_wafers["values"] = []
        tapeout = combobox_tapeout.get().strip()
        yaml_path = r"W:\50-DEVELOPMENT\TEST\Temporary data database\wafer_list.yaml"   # important to keep this list updated 
        
        if not os.path.exists(yaml_path):
            messagebox.showerror("Error", "File 'wafer_list.yaml' could not be located.")
            return

        with open(yaml_path, "r") as f:
            data = yaml.safe_load(f)
            wafers = data.get("wafers", {})
            
        filters = {"tapeout": tapeout}
        filtered = filter_wafers(filters, wafers)
        wafer_list = [k for k in filtered.keys()]
        
        root_data = r"W:\50-DEVELOPMENT\TEST\Temporary data database\Data"
        wafers_found = find_files(root_data, wafer_list, base_name, tapeout)[1]
        combobox_wafers["values"] = wafers_found
        if wafers_found:
            combobox_wafers.current(0)
        

    def validate():
        saving_file = entry_folder.get().strip()
        tapeout = combobox_tapeout.get().strip()
        base_name = entry_base.get().strip()
        generate_pptx = var_pptx.get()
        
        display_plots = var_display.get()
        save_plots = var_save.get()
        
        if not display_plots and not save_plots and not generate_pptx:
            messagebox.showerror("Error", "You must choose at least one option between Preview, Saving and Powerpoint report")
            return
            
        screen_w = root.winfo_screenwidth()
        screen_h = root.winfo_screenheight()
        
        selected_bin_str = combobox_binning.get()
        try:
            final_binning = int(selected_bin_str.split(" - ")[0])
        except (ValueError, IndexError):
            messagebox.showerror("Error", "Invalid Final Binning selection format.")
            return
        
        channel = entry_channel.get().strip()
        selected_wafer = combobox_wafers.get().strip()

        
            

        plot_choices = {
            "plot1": var_p1.get(),
            "plot2": var_p2.get(),
            "plot3": var_p3.get()
        }
        
        yaml_path = r"W:\50-DEVELOPMENT\TEST\Temporary data database\wafer_list.yaml"
        if not os.path.exists(yaml_path):
            messagebox.showerror("Error", "File 'wafer_list.yaml' not found.")
            return

        with open(yaml_path, "r") as f:
            data = yaml.safe_load(f)
            wafers = data.get("wafers", {})
            
        filters = {"tapeout": tapeout}
        filtered = filter_wafers(filters, wafers)
        wafer_list = [k for k in filtered.keys()]

        if not saving_file or not tapeout or not base_name or (tapeout == "TO5" and not channel):
            messagebox.showerror("Error", "Please fill in all standard configuration fields correctly.")
            return
        
        btn_run.config(state="disabled")
        progress.grid(row=11, column=1, pady=10)
        progress.start()

        root.after(100, lambda: start_processing(
            base_name, final_binning, channel, tapeout, saving_file, wafer_list, selected_wafer, 
            generate_pptx, screen_w, screen_h, plot_choices, save_plots, display_plots
        ))
        

    def finish_processing(save1, save2, save3, save4, display_plots):
        progress.stop()
        progress.grid_forget()
        btn_run.config(state="normal")
        messagebox.showinfo("Success", "Analysis completed successfully!")
        
        if display_plots:
            if save1: display_image(root, save1)
            if save2: display_image(root, save2)
            if save3: display_image(root, save3)
            if save4: display_image(root, save4)


    # ---- Main UI Windows Core Configuration Layout ----
    root = tk.Tk()
    root.title("Packaging General Plot Tool")
    root.geometry("900x780")
    root.configure(padx=15, pady=15)
    
    style = ttk.Style()
    style.theme_use('clam')

    desc_text = (
        "Application Description:\n"
        "For a given Tapeout name, this script extracts metrics to plot the exact number of occurrences for\n "
        "your chosen 'Final Binning'. Wafers are automatically discovered via the network (W: drive).\n"
        "All output figures are automatically exported and saved directly into your chosen destination folder below."
    )
    lbl_desc = tk.Label(root, text=desc_text, justify="left", fg="#333333", bg="#f0f0f0", bd=1, relief="solid", padx=10, pady=8, font=("Arial", 9))
    lbl_desc.grid(row=0, column=0, columnspan=3, pady=(0, 10), sticky="ew")

    tk.Label(root, text="Saving Folder:", font=("Arial", 10, "bold")).grid(row=1, column=0, padx=5, pady=5, sticky="w")
    entry_folder = tk.Entry(root, width=50, font=("Arial", 10))
    entry_folder.grid(row=1, column=1, padx=5, pady=5, sticky="w")
    tk.Button(root, text="Browse...", command=choose_folder, font=("Arial", 9)).grid(row=1, column=2, padx=5, pady=5, sticky="w")

    tk.Label(root, text="Tapeout Name:", font=("Arial", 10, "bold")).grid(row=2, column=0, padx=5, pady=5, sticky="w")
    combobox_tapeout = ttk.Combobox(root, values=["TO4.1", "TO5"], width=47, state="readonly", font=("Arial", 10))
    combobox_tapeout.set("TO5")
    combobox_tapeout.grid(row=2, column=1, padx=5, pady=5, sticky="w")

    tk.Label(root, text="Base Name:", font=("Arial", 10, "bold")).grid(row=3, column=0, padx=5, pady=5, sticky="w")
    entry_base = ttk.Combobox(root, values=["BASE_LONG", "BASE_SHORT", "ENH_LONG", "ENH_LONG_2x8"], width=50, font=("Arial", 10))
    entry_base.set("BASE_LONG")
    entry_base.grid(row=3, column=1, padx=5, pady=5, sticky="w")
    
    tk.Label(root, text="Final Binning:", font=("Arial", 10, "bold")).grid(row=4, column=0, padx=5, pady=5, sticky="w")
    combobox_binning = ttk.Combobox(root, values=BINNING_OPTIONS, width=47, state="readonly", font=("Arial", 10))
    combobox_binning.current(0)  
    combobox_binning.grid(row=4, column=1, padx=5, pady=5, sticky="w")
    
    tk.Label(root, text="Channel:", font=("Arial", 10, "bold")).grid(row=5, column=0, padx=5, pady=5, sticky="w")
    entry_channel = ttk.Combobox(root, values=["AE", "AO", "BE", "BO"], width=50, font=("Arial", 10))
    entry_channel.set("AE")
    entry_channel.grid(row=5, column=1, padx=5, pady=5, sticky="w")

    # ---- NOVEAU CADRE : SÉLECTION DES GRAPHIQUES À TRACER ----
    lf_plots = tk.LabelFrame(root, text=" Figures Selection", font=("Arial", 10, "bold"), padx=10, pady=10)
    lf_plots.grid(row=6, column=0, columnspan=3, pady=10, sticky="ew")

    var_p1 = tk.BooleanVar(value=True)
    var_p2 = tk.BooleanVar(value=True)
    var_p3 = tk.BooleanVar(value=True)

    tk.Checkbutton(lf_plots, text="Total count final binning", variable=var_p1, font=("Arial", 9)).grid(row=0, column=0, sticky="w", padx=5)
    tk.Checkbutton(lf_plots, text="Min/Max Power per die", variable=var_p2, font=("Arial", 9)).grid(row=0, column=1, sticky="w", padx=5)
    tk.Checkbutton(lf_plots, text="Wavelength of F0", variable=var_p3, font=("Arial", 9)).grid(row=0, column=2, sticky="w", padx=5)
    
    

    # Mode Options : Affichage et Sauvegarde
    var_display = tk.BooleanVar(value=True)
    var_save = tk.BooleanVar(value=True)
    
    chk_display = tk.Checkbutton(root, text="Plot figures on screen (Preview)", variable=var_display, font=("Arial", 10, "bold"), fg="#27ae60")
    chk_display.grid(row=7, column=1, pady=2, sticky="w")
    
    chk_save = tk.Checkbutton(root, text="Save figures in saving folder", variable=var_save, font=("Arial", 10, "bold"), fg="#c0392b")
    chk_save.grid(row=8, column=1, pady=2, sticky="w")

    # PowerPoint Report Checkbox
    var_pptx = tk.BooleanVar(value=False)
    chk_pptx = tk.Checkbutton(root, text="Generate PowerPoint Report (.pptx)", variable=var_pptx, font=("Arial", 10, "bold"), fg="#1f4e79")
    chk_pptx.grid(row=9, column=1, pady=2, sticky="w")

    def on_tapeout_change(event):
        if combobox_tapeout.get() in ["TO4", "TO4.1"]:
            entry_base.set("LA1")
            entry_channel.config(state="disabled")
            entry_base.config(state="disabled")
        else:
            entry_channel.config(state="normal")
            entry_channel.set("AE")
            entry_base.set("BASE_LONG")
            entry_base.config(state="normal")
            
    combobox_tapeout.bind("<<ComboboxSelected>>", on_tapeout_change)
    
    # Optional subframe area boundary zone box
    lf_optional = tk.LabelFrame(root, text=" Optional: Specific Wafer Analysis ", font=("Arial", 10, "italic"), padx=10, pady=8)
    lf_optional.grid(row=10, column=0, columnspan=3, pady=5, sticky="ew")

    tk.Button(lf_optional, text="1. Load Wafers", command=choose_wafer, font=("Arial", 9, "bold")).grid(row=0, column=0, padx=5, pady=5)
    combobox_wafers = ttk.Combobox(lf_optional, state="readonly", width=30, font=("Arial", 10))
    combobox_wafers.grid(row=0, column=1, padx=5, pady=5)
    tk.Label(lf_optional, text="(Plots the detailed internal Min-Max Power Range metrics per single Die)").grid(row=0, column=2, padx=5, pady=5)

    progress = ttk.Progressbar(root, mode="indeterminate", length=300)
    
    btn_run = tk.Button(
        root, text="Run Analysis", bg="#2ecc71", fg="white",
        activebackground="#27ae60", activeforeground="white",
        font=("Arial", 11, "bold"), bd=2, relief="groove", command=validate
    )
    btn_run.grid(row=12, column=1, pady=15, sticky="ew")

    root.mainloop()

if __name__ == "__main__":
    request_parameters()